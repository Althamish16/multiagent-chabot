"""
Advanced File Summarization Agent with LangGraph Architecture
Supports multiple file types with intelligent text extraction, chunking, and summarization
"""
import os
import json
import logging
from typing import Dict, Any, TypedDict, List, Optional, Tuple
from pathlib import Path
import tempfile
import asyncio

from langgraph.graph import StateGraph, END
from langchain_openai import AzureChatOpenAI
from langchain_text_splitters import RecursiveCharacterTextSplitter, MarkdownTextSplitter
from langchain_core.prompts import ChatPromptTemplate

from config import (
    AZURE_OPENAI_API_KEY, AZURE_OPENAI_ENDPOINT, AZURE_OPENAI_API_VERSION, AZURE_OPENAI_CHAT_DEPLOYMENT_NAME
)

# File processing libraries
try:
    import fitz  # PyMuPDF
    HAS_FITZ = True
    logging.info("✅ PyMuPDF (fitz) loaded successfully")
except ImportError as e:
    HAS_FITZ = False
    logging.warning(f"⚠️ PyMuPDF (fitz) not available: {e}")

try:
    from docx import Document
    HAS_DOCX = True
    logging.info("✅ python-docx loaded successfully")
except ImportError as e:
    HAS_DOCX = False
    logging.warning(f"⚠️ python-docx not available: {e}")

try:
    from pptx import Presentation
    HAS_PPTX = True
    logging.info("✅ python-pptx loaded successfully")
except ImportError as e:
    HAS_PPTX = False
    logging.warning(f"⚠️ python-pptx not available: {e}")

try:
    import pandas as pd
    HAS_PANDAS = True
    logging.info("✅ pandas loaded successfully")
except ImportError as e:
    HAS_PANDAS = False
    logging.warning(f"⚠️ pandas not available: {e}")


class FileSummarizerState(TypedDict):
    """State for the file summarizer LangGraph"""
    # Input
    file_content: Optional[bytes]
    file_name: str
    file_type: str
    user_request: str
    summary_mode: str  # "brief", "detailed", "executive", "technical"
    query: Optional[str]
    conversation_history: List[str]

    # Processing
    extracted_text: str
    document_structure: Dict[str, Any]
    chunks: List[Dict[str, Any]]
    chunk_summaries: List[str]

    # Output
    final_summary: str
    key_insights: List[str]
    metadata: Dict[str, Any]
    query_response: Optional[str]

    # Control
    current_step: str
    errors: List[str]
    workflow_complete: bool


class AdvancedFileSummarizerAgent:
    def __init__(self):
        self.llm = AzureChatOpenAI(
            azure_endpoint=AZURE_OPENAI_ENDPOINT,
            azure_deployment=AZURE_OPENAI_CHAT_DEPLOYMENT_NAME,
            api_version=AZURE_OPENAI_API_VERSION,
            api_key=AZURE_OPENAI_API_KEY,
            temperature=0.1
        )

        # Supported file types
        self.supported_types = {
            '.pdf': 'pdf',
            '.docx': 'docx',
            '.pptx': 'pptx',
            '.csv': 'csv',
            '.xlsx': 'xlsx',
            '.txt': 'txt',
            '.md': 'markdown',
            '.json': 'json',
            '.html': 'html'
        }

        # Text splitter for chunking
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=2000,
            chunk_overlap=200,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""]
        )

        # Build the LangGraph
        self.graph = self._build_graph()

    def _build_graph(self) -> StateGraph:
        """Build the LangGraph for file summarization"""
        workflow = StateGraph(FileSummarizerState)

        # Add nodes
        workflow.add_node("file_ingestor", self._file_ingestor)
        workflow.add_node("text_extractor", self._text_extractor)
        workflow.add_node("chunker", self._chunker)
        workflow.add_node("summarizer", self._summarizer)
        workflow.add_node("query_processor", self._query_processor)
        workflow.add_node("output_formatter", self._output_formatter)

        # Define the flow
        workflow.set_entry_point("file_ingestor")

        # Add edges
        workflow.add_edge("file_ingestor", "text_extractor")
        workflow.add_edge("text_extractor", "chunker")
        workflow.add_edge("chunker", "summarizer")

        # Conditional edges for query processing
        workflow.add_conditional_edges(
            "summarizer",
            self._should_process_query,
            {
                "query": "query_processor",
                "no_query": "output_formatter"
            }
        )

        workflow.add_edge("query_processor", "output_formatter")
        workflow.add_edge("output_formatter", END)

        return workflow.compile()

    def _should_process_query(self, state: FileSummarizerState) -> str:
        """Determine if query processing is needed"""
        return "query" if state.get("query") else "no_query"

    async def _file_ingestor(self, state: FileSummarizerState) -> FileSummarizerState:
        """Handle file ingestion and validation"""
        try:
            file_content = state.get("file_content")
            file_name = state.get("file_name", "")

            if not file_content:
                state["errors"] = ["No file content provided"]
                state["workflow_complete"] = True
                return state

            # Validate file size (50MB limit)
            if len(file_content) > 50 * 1024 * 1024:
                state["errors"] = ["File too large (max 50MB)"]
                state["workflow_complete"] = True
                return state

            # Detect file type
            file_ext = Path(file_name).suffix.lower()
            if file_ext not in self.supported_types:
                state["errors"] = [f"Unsupported file type: {file_ext}"]
                state["workflow_complete"] = True
                return state

            state["file_type"] = self.supported_types[file_ext]
            state["current_step"] = "file_ingested"

            logging.info(f"File ingested: {file_name} ({len(file_content)} bytes)")
            return state

        except Exception as e:
            state["errors"] = [f"File ingestion error: {str(e)}"]
            state["workflow_complete"] = True
            return state

    async def _text_extractor(self, state: FileSummarizerState) -> FileSummarizerState:
        """Extract text from various file formats"""
        try:
            file_content = state.get("file_content")
            file_type = state.get("file_type")
            file_name = state.get("file_name")

            if not file_content or not file_type:
                state["errors"] = ["Missing file content or type"]
                logging.error(f"Missing file content or type: content={bool(file_content)}, type={file_type}")
                return state

            logging.info(f"Starting text extraction for {file_type} file: {file_name}")

            # Create temporary file for processing
            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_type}") as temp_file:
                temp_file.write(file_content)
                temp_path = temp_file.name
                logging.info(f"Created temp file: {temp_path}")

            try:
                extracted_text = ""
                document_structure = {}

                if file_type == "pdf" and HAS_FITZ:
                    logging.info("Using PDF extractor")
                    extracted_text, document_structure = self._extract_pdf(temp_path)
                elif file_type == "docx" and HAS_DOCX:
                    logging.info("Using DOCX extractor")
                    extracted_text, document_structure = self._extract_docx(temp_path)
                elif file_type == "pptx" and HAS_PPTX:
                    logging.info("Using PPTX extractor")
                    extracted_text, document_structure = self._extract_pptx(temp_path)
                elif file_type in ["csv", "xlsx"] and HAS_PANDAS:
                    logging.info("Using spreadsheet extractor")
                    extracted_text, document_structure = self._extract_spreadsheet(temp_path, file_type)
                elif file_type == "txt":
                    logging.info("Using text extractor")
                    extracted_text, document_structure = self._extract_text(temp_path)
                elif file_type == "markdown":
                    logging.info("Using markdown extractor")
                    extracted_text, document_structure = self._extract_markdown(temp_path)
                elif file_type == "json":
                    logging.info("Using JSON extractor")
                    extracted_text, document_structure = self._extract_json(temp_path)
                elif file_type == "html":
                    logging.info("Using HTML extractor")
                    extracted_text, document_structure = self._extract_html(temp_path)
                else:
                    error_msg = f"No extractor available for {file_type} (HAS_DOCX={HAS_DOCX}, HAS_FITZ={HAS_FITZ}, HAS_PPTX={HAS_PPTX}, HAS_PANDAS={HAS_PANDAS})"
                    state["errors"] = [error_msg]
                    logging.error(error_msg)
                    return state

                if not extracted_text:
                    logging.warning(f"No text extracted from {file_name}. Document structure: {document_structure}")
                else:
                    logging.info(f"Text extracted: {len(extracted_text)} characters from {file_name}")

                state["extracted_text"] = extracted_text
                state["document_structure"] = document_structure
                state["current_step"] = "text_extracted"

                return state

            finally:
                # Clean up temp file
                try:
                    os.unlink(temp_path)
                    logging.info(f"Cleaned up temp file: {temp_path}")
                except Exception as cleanup_error:
                    logging.warning(f"Could not delete temp file {temp_path}: {cleanup_error}")

        except Exception as e:
            error_msg = f"Text extraction error: {str(e)}"
            state["errors"] = [error_msg]
            logging.error(error_msg, exc_info=True)
            return state

    def _extract_pdf(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PDF using PyMuPDF"""
        doc = fitz.open(file_path)
        text = ""
        structure = {"pages": [], "total_pages": len(doc)}

        for page_num, page in enumerate(doc):
            page_text = page.get_text()
            text += f"\n--- Page {page_num + 1} ---\n{page_text}"
            structure["pages"].append({
                "page_number": page_num + 1,
                "text_length": len(page_text)
            })

        doc.close()
        return text.strip(), structure

    def _extract_docx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from DOCX using python-docx"""
        try:
            doc = Document(file_path)
            
            # Check document properties
            try:
                core_props = doc.core_properties
                logging.info(f"DOCX properties: title='{core_props.title}', author='{core_props.author}', word_count={core_props.word_count}")
            except Exception as e:
                logging.warning(f"Could not read document properties: {e}")
            
            text = ""
            structure = {"paragraphs": [], "tables": [], "total_paragraphs": len(doc.paragraphs), "total_tables": len(doc.tables)}

            logging.info(f"DOCX extraction: {len(doc.paragraphs)} paragraphs, {len(doc.tables)} tables")

            # Extract from paragraphs
            para_count = 0
            for para_num, para in enumerate(doc.paragraphs):
                if para.text.strip():
                    text += para.text + "\n"
                    para_count += 1
                    structure["paragraphs"].append({
                        "paragraph_number": para_num + 1,
                        "text_length": len(para.text)
                    })

            logging.info(f"DOCX extraction: {para_count} non-empty paragraphs found")

            # Extract from tables
            table_count = 0
            for table_num, table in enumerate(doc.tables):
                table_text = ""
                for row_num, row in enumerate(table.rows):
                    for cell_num, cell in enumerate(row.cells):
                        for para in cell.paragraphs:
                            if para.text.strip():
                                table_text += para.text + " "
                    if table_text.strip():
                        table_text += "\n"
                
                if table_text.strip():
                    text += f"\n--- Table {table_num + 1} ---\n{table_text}"
                    table_count += 1
                    structure["tables"].append({
                        "table_number": table_num + 1,
                        "text_length": len(table_text)
                    })

            logging.info(f"DOCX extraction: {table_count} tables with content found")
            logging.info(f"DOCX extraction: Total extracted text length: {len(text)}")

            # Extract from headers and footers
            try:
                for section in doc.sections:
                    # Header
                    if section.header:
                        for para in section.header.paragraphs:
                            if para.text.strip():
                                text += f"\n--- Header ---\n{para.text}\n"
                    
                    # Footer  
                    if section.footer:
                        for para in section.footer.paragraphs:
                            if para.text.strip():
                                text += f"\n--- Footer ---\n{para.text}\n"
            except Exception as e:
                logging.warning(f"Could not extract headers/footers: {e}")

            logging.info(f"DOCX extraction: Final text length: {len(text)}")

            return text.strip(), structure
            
        except Exception as e:
            logging.error(f"Error extracting DOCX: {str(e)}")
            return "", {"error": str(e)}

    def _extract_pptx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PPTX using python-pptx"""
        prs = Presentation(file_path)
        text = ""
        structure = {"slides": [], "total_slides": len(prs.slides)}

        for slide_num, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text += shape.text + "\n"

            text += f"\n--- Slide {slide_num + 1} ---\n{slide_text}"
            structure["slides"].append({
                "slide_number": slide_num + 1,
                "text_length": len(slide_text)
            })

        return text.strip(), structure

    def _extract_spreadsheet(self, file_path: str, file_type: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from CSV/XLSX using pandas"""
        if file_type == "csv":
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)

        text = f"Columns: {', '.join(df.columns.tolist())}\n\n"
        text += df.to_string(index=False)

        structure = {
            "columns": df.columns.tolist(),
            "rows": len(df),
            "data_types": df.dtypes.to_dict()
        }

        return text, structure

    def _extract_text(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from plain text files"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()

        structure = {
            "lines": len(text.split('\n')),
            "characters": len(text)
        }

        return text, structure

    def _extract_markdown(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from markdown files"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()

        # Basic structure detection
        lines = text.split('\n')
        headers = [line for line in lines if line.startswith('#')]

        structure = {
            "headers": headers,
            "lines": len(lines),
            "characters": len(text)
        }

        return text, structure

    def _extract_json(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from JSON files"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        text = json.dumps(data, indent=2)

        structure = {
            "keys": list(data.keys()) if isinstance(data, dict) else [],
            "is_array": isinstance(data, list),
            "size": len(data) if hasattr(data, '__len__') else 0
        }

        return text, structure

    def _extract_html(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from HTML files (basic)"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()

        # Very basic HTML text extraction (remove tags)
        import re
        clean_text = re.sub(r'<[^>]+>', '', text)
        clean_text = re.sub(r'\s+', ' ', clean_text).strip()

        structure = {
            "original_length": len(text),
            "cleaned_length": len(clean_text)
        }

        return clean_text, structure

    async def _chunker(self, state: FileSummarizerState) -> FileSummarizerState:
        """Split extracted text into chunks"""
        try:
            extracted_text = state.get("extracted_text", "")
            document_structure = state.get("document_structure", {})

            logging.info(f"Chunker: extracted_text length = {len(extracted_text)}")

            if not extracted_text:
                error_msg = "No text to chunk - extracted_text is empty"
                state["errors"] = [error_msg]
                logging.error(error_msg)
                return state

            # Split text into chunks
            texts = self.text_splitter.split_text(extracted_text)
            chunks = []

            logging.info(f"Chunker: split into {len(texts)} chunks")

            for i, chunk_text in enumerate(texts):
                chunk_metadata = {
                    "chunk_id": i,
                    "text": chunk_text,
                    "length": len(chunk_text),
                    "start_char": extracted_text.find(chunk_text),
                    "end_char": extracted_text.find(chunk_text) + len(chunk_text)
                }

                # Add document structure info if available
                if "pages" in document_structure:
                    # Estimate page number based on character position
                    total_chars = len(extracted_text)
                    chars_per_page = total_chars / document_structure.get("total_pages", 1)
                    estimated_page = int(chunk_metadata["start_char"] / chars_per_page) + 1
                    chunk_metadata["estimated_page"] = estimated_page

                chunks.append(chunk_metadata)

            state["chunks"] = chunks
            state["current_step"] = "text_chunked"

            logging.info(f"Text chunked: {len(chunks)} chunks created")
            return state

        except Exception as e:
            error_msg = f"Chunking error: {str(e)}"
            state["errors"] = [error_msg]
            logging.error(error_msg, exc_info=True)
            return state

    async def _summarizer(self, state: FileSummarizerState) -> FileSummarizerState:
        """Generate summaries at chunk and document level"""
        try:
            chunks = state.get("chunks", [])
            summary_mode = state.get("summary_mode", "detailed")
            user_request = state.get("user_request", "")

            logging.info(f"Summarizer: Processing {len(chunks)} chunks in {summary_mode} mode")

            if not chunks:
                error_msg = "No chunks to summarize"
                state["errors"] = [error_msg]
                logging.error(error_msg)
                return state

            # Summarize each chunk
            chunk_summaries = []
            for i, chunk in enumerate(chunks):
                logging.info(f"Summarizing chunk {i+1}/{len(chunks)} ({len(chunk['text'])} chars)")
                summary = await self._summarize_chunk(chunk["text"], summary_mode)
                chunk_summaries.append(summary)
                logging.info(f"Chunk {i+1} summary: {len(summary)} chars")

            state["chunk_summaries"] = chunk_summaries
            logging.info(f"All chunk summaries complete: {len(chunk_summaries)} summaries")

            # Generate document-level summary
            conversation_history = state.get("conversation_history", [])
            logging.info(f"Generating document-level summary with {len(conversation_history)} history items")
            
            document_summary = await self._summarize_document(
                chunk_summaries, summary_mode, user_request, conversation_history
            )
            
            logging.info(f"Document summary generated: {len(document_summary)} chars")

            # Extract key insights
            logging.info("Extracting key insights...")
            key_insights = await self._extract_key_insights(document_summary)
            logging.info(f"Key insights extracted: {len(key_insights)} insights")

            state["final_summary"] = document_summary
            state["key_insights"] = key_insights
            state["current_step"] = "summaries_generated"

            logging.info(f"Summaries generated: {len(chunk_summaries)} chunks, {len(key_insights)} insights")
            return state

        except Exception as e:
            error_msg = f"Summarization error: {str(e)}"
            state["errors"] = [error_msg]
            logging.error(error_msg, exc_info=True)
            return state

    async def _summarize_chunk(self, chunk_text: str, summary_mode: str) -> str:
        """Summarize a single chunk"""
        mode_prompts = {
            "brief": "Summarize this text in 1-2 sentences:",
            "detailed": "Provide a detailed summary of this text:",
            "executive": "Provide an executive summary highlighting key business implications:",
            "technical": "Provide a technical summary with key details and specifications:"
        }

        prompt = ChatPromptTemplate.from_template(f"""
        {mode_prompts.get(summary_mode, mode_prompts['detailed'])}

        Text: {{chunk_text}}

        Focus on the most important information and maintain factual accuracy.
        """)

        chain = prompt | self.llm
        response = await chain.ainvoke({"chunk_text": chunk_text})
        return response.content.strip()

    async def _summarize_document(self, chunk_summaries: List[str], summary_mode: str, user_request: str, conversation_history: List[str] = None) -> str:
        """Generate document-level summary from chunk summaries"""
        combined_summaries = "\n\n".join(chunk_summaries)

        mode_instructions = {
            "brief": "Create a brief 2-3 sentence summary of the entire document.",
            "detailed": "Create a comprehensive summary covering all main points and details.",
            "executive": "Create an executive summary focusing on key decisions, outcomes, and business impact.",
            "technical": "Create a technical summary with specifications, methodologies, and detailed findings."
        }

        prompt = ChatPromptTemplate.from_template(f"""
        You are creating a {summary_mode} summary of a document based on chunk summaries.

        User context: {user_request}
        
        Recent conversation history (for additional context):
        {chr(10).join(conversation_history[-5:]) if conversation_history else "No recent conversation"}

        {mode_instructions.get(summary_mode, mode_instructions['detailed'])}

        Chunk summaries:
        {{combined_summaries}}

        Generate a cohesive {summary_mode} summary that captures the essence of the entire document.
        Consider the conversation history to provide contextually relevant information.
        """)

        chain = prompt | self.llm
        response = await chain.ainvoke({"combined_summaries": combined_summaries})
        return response.content.strip()

    async def _extract_key_insights(self, document_summary: str) -> List[str]:
        """Extract key insights from the document summary"""
        prompt = ChatPromptTemplate.from_template("""
        Extract 3-5 key insights from this document summary. Each insight should be:
        - Concise (1 sentence)
        - Actionable or informative
        - Factual

        Summary: {document_summary}

        Return as a JSON array of strings.
        """)

        chain = prompt | self.llm
        response = await chain.ainvoke({"document_summary": document_summary})

        try:
            insights = json.loads(response.content.strip())
            return insights if isinstance(insights, list) else []
        except:
            # Fallback: split by newlines
            return [line.strip() for line in response.content.split('\n') if line.strip()][:5]

    async def _query_processor(self, state: FileSummarizerState) -> FileSummarizerState:
        """Process user queries about the document"""
        try:
            query = state.get("query")
            chunks = state.get("chunks", [])
            chunk_summaries = state.get("chunk_summaries", [])

            if not query or not chunks:
                state["query_response"] = "Unable to process query"
                return state

            # Simple relevance matching (could be enhanced with embeddings)
            relevant_chunks = []
            query_lower = query.lower()

            for i, chunk in enumerate(chunks):
                chunk_text_lower = chunk["text"].lower()
                if any(word in chunk_text_lower for word in query_lower.split()):
                    relevant_chunks.append({
                        "chunk_id": i,
                        "text": chunk["text"],
                        "summary": chunk_summaries[i] if i < len(chunk_summaries) else ""
                    })

            # Generate query response
            context_text = "\n\n".join([f"Section {i+1}: {chunk['text'][:500]}..." for i, chunk in enumerate(relevant_chunks[:3])])

            prompt = ChatPromptTemplate.from_template("""
            Answer this query about the document using the provided context:

            Query: {query}

            Context from relevant sections:
            {context_text}

            Provide a direct, factual answer based on the context. If the context doesn't contain enough information, say so.
            """)

            chain = prompt | self.llm
            response = await chain.ainvoke({
                "query": query,
                "context_text": context_text
            })

            state["query_response"] = response.content.strip()
            return state

        except Exception as e:
            state["query_response"] = f"Query processing error: {str(e)}"
            return state

    async def _output_formatter(self, state: FileSummarizerState) -> FileSummarizerState:
        """Format the final output with metadata"""
        try:
            final_summary = state.get("final_summary", "")
            key_insights = state.get("key_insights", [])
            extracted_text = state.get("extracted_text", "")
            file_name = state.get("file_name", "")
            summary_mode = state.get("summary_mode", "detailed")
            query_response = state.get("query_response")

            # Calculate metadata
            original_length = len(extracted_text)
            summary_length = len(final_summary)
            reduction_percentage = ((original_length - summary_length) / original_length * 100) if original_length > 0 else 0

            metadata = {
                "file_name": file_name,
                "summary_mode": summary_mode,
                "original_text_length": original_length,
                "summary_length": summary_length,
                "reduction_percentage": round(reduction_percentage, 1),
                "num_chunks": len(state.get("chunks", [])),
                "num_insights": len(key_insights),
                "processing_timestamp": str(asyncio.get_event_loop().time()),
                "agent_version": "1.0.0"
            }

            state["metadata"] = metadata
            state["workflow_complete"] = True
            state["current_step"] = "output_formatted"

            logging.info(f"Output formatted: {summary_length} chars summary, {reduction_percentage:.1f}% reduction")
            return state

        except Exception as e:
            state["errors"] = [f"Output formatting error: {str(e)}"]
            return state

    async def _check_cached_summary(self, file_name: str, user_request: str, conversation_history: List[str]) -> Optional[Dict[str, Any]]:
        """Check if we can use a cached summary instead of re-processing the document"""
        if not conversation_history:
            return None
            
        try:
            # Look for previous document summaries in conversation history
            cached_summaries = []
            for msg in conversation_history:
                if msg.startswith("Assistant:"):
                    # Extract the response content
                    response_content = msg.replace("Assistant: ", "")
                    # Look for substantial responses that could be document summaries
                    # Check for length (>500 chars), and contains document-related terms
                    if (len(response_content) > 500 and 
                        any(term in response_content.lower() for term in ["summary", "document", "analysis", "key insights", "extracted"])):
                        cached_summaries.append(response_content)
            
            if not cached_summaries:
                return None
                
            # Use the most recent substantial summary
            combined_summary = cached_summaries[-1]
            
            # Determine if this is a question that can be answered from the summary
            question_keywords = ["what", "how", "why", "when", "where", "who", "which", "can you", "tell me", "explain", "?"]
            is_question = any(keyword in user_request.lower() for keyword in question_keywords) or "?" in user_request
            
            if is_question:
                # Use LLM to answer the question based on cached summary
                answer = await self._answer_from_cached_summary(user_request, combined_summary, file_name)
                if answer:
                    logging.info(f"Answered question from cached summary for {file_name}")
                    return {
                        "status": "success",
                        "query_response": answer,
                        "summary": combined_summary,
                        "key_insights": [],
                        "metadata": {
                            "file_name": file_name,
                            "cached": True,
                            "answer_type": "from_cached_summary"
                        },
                        "file_type": "cached",
                        "processing_steps": "cached_response"
                    }
            else:
                # For non-questions, return the cached summary directly
                logging.info(f"Returning cached summary for {file_name}")
                return {
                    "status": "success",
                    "summary": combined_summary,
                    "key_insights": [],
                    "metadata": {
                        "file_name": file_name,
                        "cached": True,
                        "summary_type": "from_cache"
                    },
                    "file_type": "cached",
                    "processing_steps": "cached_summary"
                }
                
        except Exception as e:
            logging.warning(f"Error checking cached summary: {e}")
            return None
            
        return None

    async def _answer_from_cached_summary(self, question: str, cached_summary: str, file_name: str) -> Optional[str]:
        """Generate an answer to a question using the cached document summary"""
        try:
            prompt = ChatPromptTemplate.from_template("""
            Answer this question about the document "{file_name}" using the provided summary.
            If the summary doesn't contain enough information to answer the question, say so clearly.

            Question: {question}

            Document Summary:
            {cached_summary}

            Provide a direct, factual answer based on the summary. Keep your response concise.
            """)

            chain = prompt | self.llm
            response = await chain.ainvoke({
                "question": question,
                "cached_summary": cached_summary,
                "file_name": file_name
            })
            
            answer = response.content.strip()
            
            # Check if the answer indicates insufficient information
            insufficient_indicators = ["doesn't contain", "not enough information", "cannot answer", "insufficient"]
            if any(indicator in answer.lower() for indicator in insufficient_indicators):
                return None  # Force re-processing
                
            return answer
            
        except Exception as e:
            logging.warning(f"Error generating answer from cached summary: {e}")
            return None

    async def process_file(self,
                          file_content: bytes,
                          file_name: str,
                          user_request: str = "",
                          summary_mode: str = "detailed",
                          query: Optional[str] = None,
                          conversation_history: Optional[List[str]] = None) -> Dict[str, Any]:
        """Main entry point for file processing"""
        logging.info(f"Processing file: {file_name} ({len(file_content)} bytes)")

        try:
            conversation_history = conversation_history or []
            
            # Check if document has been processed before and if we can use cached summary
            cached_response = await self._check_cached_summary(file_name, user_request, conversation_history)
            if cached_response:
                logging.info(f"Using cached summary for {file_name}")
                return cached_response
            
            # Initialize state
            initial_state: FileSummarizerState = {
                "file_content": file_content,
                "file_name": file_name,
                "file_type": "",
                "user_request": user_request,
                "summary_mode": summary_mode,
                "query": query,
                "conversation_history": conversation_history,
                "extracted_text": "",
                "document_structure": {},
                "chunks": [],
                "chunk_summaries": [],
                "final_summary": "",
                "key_insights": [],
                "metadata": {},
                "query_response": None,
                "current_step": "",
                "errors": [],
                "workflow_complete": False
            }

            # Execute the graph
            final_state = await self.graph.ainvoke(initial_state)

            # Check for errors
            if final_state.get("errors"):
                return {
                    "status": "error",
                    "message": f"Processing failed: {'; '.join(final_state['errors'])}",
                    "errors": final_state["errors"]
                }

            # Format response
            response = {
                "status": "success",
                "summary": final_state.get("final_summary", ""),
                "key_insights": final_state.get("key_insights", []),
                "metadata": final_state.get("metadata", {}),
                "file_type": final_state.get("file_type", ""),
                "processing_steps": final_state.get("current_step", "")
            }

            if final_state.get("query_response"):
                response["query_response"] = final_state["query_response"]

            return response

        except Exception as e:
            logging.error(f"File processing error: {str(e)}")
            return {
                "status": "error",
                "message": f"Unexpected error: {str(e)}",
                "errors": [str(e)]
            }